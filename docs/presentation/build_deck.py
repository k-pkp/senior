"""Rebuild the 24 August meeting deck.

Keeps every figure the original used and its visual language -- white ground,
Canva Sans, a title block at the top left -- and changes three things: every
slide states its finding in words as well as in a plot, the figures are placed
on one grid instead of wherever they landed, and the Stage 4 conclusion is
corrected. The original argued for alpha shape; the pipeline now ships Poisson
with alpha shape as a fallback, because the defect that argument rested on was
in Stage 5.
"""
import copy, os, sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

SOURCE, TARGET = sys.argv[1], sys.argv[2]
FIGURES = sys.argv[3]

BOLD, BODY = "Canva Sans Bold", "Canva Sans"
INK = RGBColor(0x00, 0x00, 0x00)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
WARN = RGBColor(0xC0, 0x63, 0x00)
GOOD = RGBColor(0x0A, 0x6E, 0x2E)
RULE = RGBColor(0xD8, 0xD8, 0xD8)

SLIDE_W, SLIDE_H = 18288000, 10287000
MARGIN = 1028700
COLUMN_W = 4800600                      # left text column
GUTTER = 685800
CONTENT_X = MARGIN + COLUMN_W + GUTTER
CONTENT_W = SLIDE_W - MARGIN - CONTENT_X
TITLE_Y = 800100
BODY_Y = 2400300
CONTENT_Y = 1600200
CONTENT_H = SLIDE_H - CONTENT_Y - 800100


def textbox(slide, x, y, w, h, blocks, align=PP_ALIGN.LEFT):
    """`blocks` is a list of (text, size_pt, bold, colour, space_before_pt)."""
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    frame = box.text_frame
    frame.word_wrap = True
    for index, (text, size, bold, colour, space) in enumerate(blocks):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.space_before = Pt(space)
        para.line_spacing = 1.15
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = BOLD if bold else BODY
        run.font.color.rgb = colour
    return box


def heading(slide, eyebrow, title):
    blocks = []
    if eyebrow:
        blocks.append((eyebrow, 18, True, MUTED, 0))
    blocks.append((title, 36, True, INK, 8 if eyebrow else 0))
    return textbox(slide, MARGIN, TITLE_Y, SLIDE_W - 2 * MARGIN, 1400000, blocks)


def bullets(slide, lines, x=MARGIN, y=BODY_Y, w=COLUMN_W, size=19):
    blocks = []
    for index, line in enumerate(lines):
        text, colour = line if isinstance(line, tuple) else (line, INK)
        blocks.append(("•  " + text, size, False, colour, 0 if index == 0 else 13))
    return textbox(slide, x, y, w, SLIDE_H - y - 700000, blocks)


def footnote(slide, text, colour=MUTED, y=8686800):
    """A rule and one line across the foot of the slide -- the 'so what'."""
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(MARGIN), Emu(y),
                                  Emu(SLIDE_W - 2 * MARGIN), Emu(12700))
    line.fill.solid(); line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False
    textbox(slide, MARGIN, y + 152400, SLIDE_W - 2 * MARGIN, 900000,
            [(text, 21, False, colour, 0)])


def fit(path, x, y, w, h):
    """Largest box with the image's aspect ratio that fits, centred in (x,y,w,h)."""
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    dw, dh = int(iw * scale), int(ih * scale)
    return int(x + (w - dw) / 2), int(y + (h - dh) / 2), dw, dh


def picture(slide, path, x, y, w, h):
    px, py, pw, ph = fit(path, x, y, w, h)
    return slide.shapes.add_picture(path, Emu(px), Emu(py), Emu(pw), Emu(ph))


def table(slide, rows, x, y, w, col_widths=None, size=16, header=True):
    row_h = int(size * 12700 * 2.6)
    shape = slide.shapes.add_table(len(rows), len(rows[0]),
                                   Emu(x), Emu(y), Emu(w),
                                   Emu(len(rows) * row_h))

    grid = shape.table
    grid.first_row = header
    grid.horz_banding = False
    if col_widths:
        total = sum(col_widths)
        for index, share in enumerate(col_widths):
            grid.columns[index].width = Emu(int(w * share / total))
    for row_obj in grid.rows:
        row_obj.height = Emu(row_h)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = grid.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Emu(91440)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(value)
            run.font.size = Pt(size)
            run.font.bold = (r == 0 and header)
            run.font.name = BOLD if (r == 0 and header) else BODY
            run.font.color.rgb = INK
    return shape


presentation = Presentation(SOURCE)
blank = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 \
        else presentation.slide_layouts[0]


def find(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def retitle(slide, eyebrow, title):
    """Replace the Canva title/subtitle pair with one heading block."""
    for name in ("TextBox 3", "TextBox 4"):
        shape = find(slide, name)
        if shape is not None:
            shape._element.getparent().remove(shape._element)
    heading(slide, eyebrow, title)


def place_figure(slide, x=CONTENT_X, y=CONTENT_Y, w=CONTENT_W, h=CONTENT_H):
    """Move the Canva picture-filled freeform onto the content grid."""
    shape = find(slide, "Freeform 2")
    if shape is None:
        return
    aspect = shape.width / shape.height
    if aspect >= w / h:
        dw = w; dh = int(w / aspect)
    else:
        dh = h; dw = int(h * aspect)
    shape.left, shape.top = Emu(int(x + (w - dw) / 2)), Emu(int(y + (h - dh) / 2))
    shape.width, shape.height = Emu(dw), Emu(dh)


slides = list(presentation.slides)

# ── 1 · title ────────────────────────────────────────────────────────────
title_slide = slides[0]
for shape in list(title_slide.shapes):
    shape._element.getparent().remove(shape._element)
textbox(title_slide, MARGIN, 3800000, SLIDE_W - 2 * MARGIN, 3000000, [
    ("Measuring limb volume from six phone photos", 60, True, INK, 0),
    ("A reworked VGGT pipeline — what changed, what it is worth, "
     "and one conclusion that was wrong", 26, False, MUTED, 22),
    ("Senior project · progress meeting · 24 August 2026", 20, False, MUTED, 26),
])

# ── 2 · what the system does (new, moved into second position) ──────────
context = presentation.slides.add_slide(blank)
context_id = presentation.slides._sldIdLst[-1]
heading(context, "THE PROBLEM", "Six phone photos in, a volume in cm³ out")
bullets(context, [
    "Input: 6 to 12 photographs of a limb standing beside a 14 cm ArUco cube.",
    "VGGT reconstructs the scene as a point cloud. The cube is the only thing in "
    "the frame whose size is known, so it supplies the real-world scale — "
    "2744 cm³ by construction, 14³.",
    "The limb is segmented, the ghost surface removed, the cut plane found at the "
    "marker band, and the survivor turned into a closed mesh.",
    "Output: the volume that mesh encloses, integrated exactly — and a statement "
    "of whether it is one closed solid, so the number can be trusted or refused.",
], x=MARGIN, y=2590800, w=SLIDE_W - 2 * MARGIN, size=22)
footnote(context, "Seven stages. Two of them stop and wait for a person: the "
                  "framing verdict, and the cut plane.", INK, y=8153400)

# ── 3 · before ───────────────────────────────────────────────────────────
retitle(slides[1], "WHERE IT STARTED",
        "Before — six stages, nothing checked the input")
place_figure(slides[1], MARGIN, 4400000, SLIDE_W - 2 * MARGIN, 1800000)
bullets(slides[1], [
    "The photographs went straight to VGGT, which centre-crops them itself.",
    "One reconstruction method, one fixed set of thresholds, no way to intervene.",
    "The reference mesh never closed, so its volume was estimated by flooding a "
    "voxel grid — an approximation with no error signal when it leaks.",
], x=MARGIN, y=2300000, w=SLIDE_W - 2 * MARGIN, size=20)
footnote(slides[1], "One pass, start to finish, with nothing between the "
                    "photographs and a number.")

# ── 4 · after ────────────────────────────────────────────────────────────
retitle(slides[2], "WHERE IT IS NOW",
        "After — seven stages, two can stop the run")
place_figure(slides[2], MARGIN, 4400000, SLIDE_W - 2 * MARGIN, 1800000)
bullets(slides[2], [
    "Stage 0 is new: it judges the capture before a GPU is touched.",
    "Stage 3 measures the marker's colour from the capture instead of assuming it, "
    "and the cut waits for a human to confirm the plane.",
    "Stage 4 chooses its method and Stage 5 proves the result is one closed solid — "
    "if it is not, Stage 4 is re-run with a method that cannot fail to close.",
], x=MARGIN, y=2300000, w=SLIDE_W - 2 * MARGIN, size=20)
footnote(slides[2], "Every stage writes its own directory, so any one of them "
                    "can be re-run without repeating the rest.")

# ── 5 · stage 0 ──────────────────────────────────────────────────────────
retitle(slides[3], "STAGE 0 · FRAMING GATE",
        "A capture can now be refused before any of it is measured")
place_figure(slides[3])
bullets(slides[3], [
    "VGGT crops a 9:16 photo to a square and discards 43.8% of it, with no "
    "regard for where the reference cube is.",
    "On IMG_4462 that cuts 16% off the cube — and the cube sets the scale for "
    "every number the run reports. A truncated cube still reconstructs; it just "
    "reconstructs smaller.",
    "Stage 0 places its own window: 100% of the cube on all five frames it can crop.",
    ("Three verdicts. IMG_4458 is a warning, not a rejection — a square holding "
     "the cube exists, but none holds the cube and the band together, so the "
     "frame is measured through VGGT's crop and the report says so.", WARN),
])

# ── 6 · stage 3 · marker colour ──────────────────────────────────────────
retitle(slides[4], "STAGE 3 · WHERE TO CUT",
        "The marker's colour is measured, not assumed")
place_figure(slides[4])
bullets(slides[4], [
    "main used a hardcoded khaki window that never looked at the photographs.",
    "Tracing only the cord's darkest pixel finds 45 points, and the cut it fits "
    "tilts 27.1° from vertical.",
    "Dilating the trace by ±3 px finds 296 — 6.6× more evidence — and the cut "
    "tilts 19.0°, matching the limb's own 19.0° lean.",
    ("The cut plane sets where the limb ends. A plane fitted to 45 points is the "
     "difference between measuring the calf and measuring part of the ankle.", MUTED),
])

# ── 7 · stage 3 · ghost ──────────────────────────────────────────────────
retitle(slides[5], "STAGE 3 · THE GHOST",
        "VGGT emits the same surface more than once")
place_figure(slides[5])
bullets(slides[5], [
    "Each camera group registers the limb slightly differently, so the shell "
    "arrives about 2.7 mm thick instead of thin.",
    "Voxel downsampling collapses 6.4 points per voxel into one: 114,282 → 17,979 "
    "points. That ratio is the measurement of how duplicated the surface was.",
    "normal_aware_filter then removes only 535 stragglers — 2.9%.",
    "MLS deletes nothing at all and still halves the shell: 1.63 mm → 0.66 mm RMS.",
])

# ── new · how MLS works, before the plane-vs-quadratic argument ─────────
mls = presentation.slides.add_slide(blank)
mls_id = presentation.slides._sldIdLst[-1]
heading(mls, "STAGE 3 · SMOOTHING", "How MLS works, in four steps")
picture(mls, os.path.join(FIGURES, "mls_how_it_works.png"),
        MARGIN, 2209800, SLIDE_W - 2 * MARGIN, 6858000)
footnote(mls, "Moving Least Squares. Each point is re-fitted against its own "
              "neighbours — collect them inside a radius, take a local frame "
              "from their spread, fit height as a curved function of position, "
              "and move the point onto that surface.", INK, y=9296400)

# ── 8 · stage 3 · MLS ────────────────────────────────────────────────────
retitle(slides[6], "STAGE 3 · SMOOTHING",
        "MLS has to be quadratic")
place_figure(slides[6])
bullets(slides[6], [
    "MLS projects each point onto a surface fitted to its neighbours. Fit a plane "
    "to a curved limb and the plane sits inside it, so the outline shrinks.",
    "On this slice: plane MLS loses 1.36% of the cross-section area, quadratic "
    "MLS 0.05%.",
    "Across 40 slices the gap is +1.10 pp and positive in 100% of them — not noise.",
    "Both collapse the shell equally well. Only the quadratic keeps the shape.",
])

# ── 9 · stage 4 · alpha ladder ───────────────────────────────────────────
retitle(slides[7], "STAGE 4 · THE GUARANTEE",
        "α is chosen by closure, not by eye")
place_figure(slides[7])
bullets(slides[7], [
    "The ladder runs α from 8× to 200× the mean point spacing and takes the "
    "smallest α that is watertight and has Euler characteristic χ = 2.",
    "χ = 2 − 2g, so χ = 2 means genus 0: one closed surface with no handles. "
    "It is the only case where 'the volume enclosed' is defined at all.",
    "Below the mark the surface is riddled with holes — χ reaches −2876 — and no "
    "signed volume exists.",
    ("Because the ladder selects on χ, it cannot fail to produce χ = 2. That is "
     "exactly what makes alpha shape the fallback in the pipeline that ships.", GOOD),
])

# ── 10 · stage 4 · comparison, superseded ────────────────────────────────
retitle(slides[8], "STAGE 4 · THE COMPARISON",
        "Three methods, same points, wrong conclusion")
place_figure(slides[8])
bullets(slides[8], [
    "Ball pivoting interpolates every input point exactly — p95 distance 0.00 mm — "
    "and still reports +30.3%. Looking right is not being right.",
    "Poisson fitted the points closest (1.02 mm) but came out of the pipeline at "
    "χ = 0 with 11 components. So alpha shape was chosen.",
    ("That conclusion was wrong, and it was wrong for a reason worth showing: "
     "the defect was not in Poisson.", WARN),
])
footnote(slides[8], "The figure is exactly as it was measured. What it concludes "
                    "no longer holds, and the next slide is why.", WARN)

# ── 11 · new · stage 5 was the bug ───────────────────────────────────────
slide = presentation.slides.add_slide(blank)
heading(slide, "STAGE 5 · THE REPAIR STAGE",
        "The stage meant to fix this was doing half its job")
bullets(slide, [
    "workers/meshfix_worker.py called pymeshfix.fill_holes() and nothing else.",
    "fill_holes closes the boundary but leaves self-intersections and "
    "non-manifold edges — watertight, and still not a solid.",
    "repair() removes those too. Every Poisson mesh had been arriving at Stage 6 "
    "closed and topologically wrong, and Poisson was blamed for it.",
    ("Cost of the fix to the reported volume: 0.005% — 1069.56 → 1069.61 cm³. "
     "Alpha-shape meshes are byte-identical either way; they arrive closed, so "
     "Stage 5 is a no-op on them.", MUTED),
], w=7000000, size=20)
table(slide, [
    ["χ after Stage 5", "fill_holes", "repair()"],
    ["small_leg · reference cube", "30", "2"],
    ["small_leg · limb", "22", "2"],
    ["short_leg · reference cube", "10", "2"],
    ["short_leg · limb (uncut, with foot)", "14", "−16"],
    ["48-config sweep, χ = 2 on both", "0 of 48", "38 of 46"],
], x=8915400, y=2590800, w=8344200, col_widths=[2.2, 1, 1], size=20)
footnote(slide, "The reported volume moves 0.005% when Stage 5 is fixed — "
                "1069.56 to 1069.61 cm³. What changes is the topology, not the shape.")

# ── 12 · new · what ships ────────────────────────────────────────────────
slide = presentation.slides.add_slide(blank)
heading(slide, "STAGE 4 · WHAT SHIPS NOW",
        "Poisson everywhere, alpha shape as the fallback")
bullets(slide, [
    "Poisson fits the points 1.8× closer on small_leg and 10× closer on short_leg.",
    "It is insensitive to its own parameters: twelve depth/trim combinations span "
    "0.32% of limb volume. trim > 0 is required — every trim = 0 run gives χ = 4.",
    "What it costs, stated plainly: Poisson has no genus guarantee. On short_leg "
    "it closes at χ = −18, about ten handles, and reads 22% below the alpha answer.",
    "So Stage 5 now checks χ, and any object that does not reach 2 is reconstructed "
    "again with alpha shape. Verified firing on short_leg and not firing on small_leg.",
], w=7000000, size=20)
table(slide, [
    ["small_leg", "limb cm³", "box χ", "limb χ", "p95"],
    ["alpha / alpha", "1081.94", "2", "2", "2.39 mm"],
    ["alpha cube + Poisson limb", "1070.85", "2", "2", "1.30 mm"],
    ["Poisson / Poisson  ← ships", "1074.32", "2", "2", "1.30 mm"],
    ["short_leg · no cut", "", "", "", ""],
    ["alpha / alpha", "2763.24", "2", "2", "21.80 mm"],
    ["Poisson / Poisson", "2146.39", "2", "−18", "2.19 mm"],
], x=8915400, y=2590800, w=8344200,
   col_widths=[2.5, 1.05, 0.7, 0.7, 1.05], size=18)
footnote(slide, "short_leg is the honest counter-example: no band means no cut, "
                "so the object is the whole foot — toes, arch, the gap under the "
                "instep. Genuine topology — and the reason the fallback exists.", WARN)

# ── 13 · new · results ───────────────────────────────────────────────────
slide = presentation.slides.add_slide(blank)
heading(slide, "RESULTS", "Where the numbers stand, on inputs/small_leg")
table(slide, [
    ["", "main", "current"],
    ["wall clock", "136.7 s", "80 s"],
    ["capture validation", "none", "5 pass, 1 warn, 0 reject"],
    ["marker colour", "hardcoded khaki", "measured per capture"],
    ["reference mesh watertight", "no", "yes, χ = 2"],
    ["volume method", "warp+floodfill", "watertight — same code"],
    ["limb volume, cm³", "1073.98 (flooded)", "1074.32 (integrated)"],
], x=MARGIN, y=2590800, w=9144000, col_widths=[1.9, 1.25, 1.85], size=18)
bullets(slide, [
    ("The reference cube reports exactly 2744.00 cm³ on every run. That is an "
     "identity, not an accuracy result: scale is derived as (real / mesh)^(1/3) "
     "from that same cube, so it can never disagree with itself.", WARN),
    "The watertight row is the one that matters. main's reference mesh never "
    "closes, so its volume cannot be integrated — it is flooded, and a flood "
    "leaks through any hole without saying so.",
    "The current tree searches until the mesh closes and χ = 2, then integrates "
    "the surface exactly.",
    ("Stage 6 is main's code, byte for byte — 618 added lines, all comments. "
     "Its first tier was always the exact signed volume; main never reached it "
     "because its meshes were open. The improvement is upstream.", WARN),
    ("The two limb figures agreeing to 0.03% is not evidence of anything: one is "
     "a flood fill of a mesh with holes, the other an exact integral. They are "
     "not the same measurement.", MUTED),
], x=10515600, y=2590800, w=6743700, size=19)

# ── 14 · new · limits ────────────────────────────────────────────────────
slide = presentation.slides.add_slide(blank)
heading(slide, "WHAT IS NOT KNOWN", "The limits, stated before someone else states them")
bullets(slide, [
    "There is no ground truth for the limb. Nothing in this deck is an accuracy "
    "measurement against a known volume.",
    "est_325's 325 ml is a fill volume, not external displacement, so it does not "
    "settle the error target either.",
    "Poisson has no genus guarantee, and short_leg shows what that looks like: "
    "χ = −18 on an uncut foot, where toes and the gap under the instep are "
    "genuine topological complexity rather than a defect.",
    "Both Poisson successes are the same capture. Whether it holds on a cut limb "
    "from a second capture is untested.",
    "Stage 0's two new thresholds — a band at most 0.35 of the limb, seen in at "
    "least 0.6 of the frames — have been tested on one banded capture only.",
], x=MARGIN, y=2400300, w=SLIDE_W - 2 * MARGIN, size=21)
footnote(slide, "Next: a second banded capture, to test the fallback on a cut "
                "limb and the two Stage 0 thresholds on something other than "
                "the capture they were written for.", INK, y=8153400)

presentation.slides._sldIdLst.insert(1, context_id)   # straight after the title
presentation.slides._sldIdLst.insert(7, mls_id)       # before plane-vs-quadratic

presentation.save(TARGET)
print("wrote", TARGET, "-", len(presentation.slides.__iter__.__self__._sldIdLst), "slides")
